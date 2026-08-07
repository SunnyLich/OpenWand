"""Tests for test document context contract."""

from __future__ import annotations

import importlib
import sys
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

import pytest


@pytest.fixture()
def win_context_fetcher():
    with patch.object(sys, "platform", "win32"):
        import core.context_fetcher as context_fetcher

        module = importlib.reload(context_fetcher)
        try:
            yield module
        finally:
            importlib.reload(context_fetcher)


@pytest.mark.parametrize(
    ("title", "process_name", "open_path", "doc_name"),
    [
        (
            "Budget.xlsx - Excel",
            "EXCEL.EXE",
            r"C:\Users\TestUser\Documents\Budget.xlsx",
            "Budget.xlsx",
        ),
        (
            "Quarterly review.pptx - PowerPoint",
            "POWERPNT.EXE",
            r"C:\Users\TestUser\Documents\Quarterly review.pptx",
            "Quarterly review.pptx",
        ),
        (
            "Proposal.docx - Word Localized",
            "WINWORD.EXE",
            r"C:\Users\TestUser\Documents\Proposal.docx",
            "Proposal.docx",
        ),
        (
            "Untitled 1 \u2014 LibreOffice Calc",
            "soffice.bin",
            r"C:\Users\TestUser\Documents\Untitled 1.ods",
            "Untitled 1",
        ),
        (
            "Notes.odt \u2014 LibreOffice Writer",
            "soffice.bin",
            r"C:\Users\TestUser\Documents\Notes.odt",
            "Notes.odt",
        ),
        (
            "laptop walmart invoice - PDF-XChange Editor",
            "PXCEditor.exe",
            r"C:\Users\TestUser\Documents\laptop walmart invoice.pdf",
            "laptop walmart invoice",
        ),
        (
            "Manual.pdf - Adobe Acrobat",
            "Acrobat.exe",
            r"C:\Users\TestUser\Documents\Manual.pdf",
            "Manual.pdf",
        ),
        (
            "Receipt - SumatraPDF",
            "SumatraPDF.exe",
            r"C:\Users\TestUser\Documents\Receipt.pdf",
            "Receipt",
        ),
        (
            "Draft.md - Notepad++",
            "notepad++.exe",
            r"C:\Users\TestUser\Documents\Draft.md",
            "Draft.md",
        ),
    ],
)
def test_common_document_apps_resolve_open_file_paths(
    win_context_fetcher,
    title,
    process_name,
    open_path,
    doc_name,
):
    cf = win_context_fetcher
    win = cf.WindowInfo(title=title, process_name=process_name, pid=101)

    assert cf._extract_doc_name_from_window(win) == doc_name
    with patch.object(cf, "_win_open_files_for_pid", return_value=[open_path]), \
         patch.object(cf, "_fetch_recent_files", return_value=[]):
        assert cf._resolve_doc_path(win) == open_path


def test_unsaved_common_document_window_uses_visible_text_fallback(win_context_fetcher):
    cf = win_context_fetcher
    cf._context_window = cf.WindowInfo(
        title="Untitled 1 \u2014 LibreOffice Calc",
        process_name="soffice.bin",
        pid=101,
        hwnd=777,
    )

    with patch.object(cf, "_enumerate_open_doc_windows", return_value=[]), \
         patch.object(cf, "_get_window_text_uia", return_value="A1\tB1\nA2\tB2"):
        assert cf.get_all_open_document_window_texts() == [
            ("Untitled 1", "A1\tB1\nA2\tB2")
        ]


def test_text_and_csv_document_readers(tmp_path):
    pytest.importorskip("anydoc")
    from core.llm_clients import client as llm

    text_path = tmp_path / "notes.md"
    text_path.write_text("# Notes\nhello text", encoding="utf-8")
    csv_path = tmp_path / "table.csv"
    csv_path.write_text("name,value\nalpha,1", encoding="utf-8")

    assert "hello text" in llm.read_document_file(str(text_path))
    csv_text = llm.read_document_file(str(csv_path))
    assert "| name | value |" in csv_text
    assert "| alpha | 1 |" in csv_text


def test_anydoc_docx_reader_preserves_heading_and_table_structure(tmp_path):
    pytest.importorskip("anydoc")
    pytest.importorskip("docx")
    from docx import Document

    from core.llm_clients import client as llm

    path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_heading("Contract heading", level=1)
    doc.add_paragraph("Docx contract text")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Item"
    table.cell(0, 1).text = "Total"
    table.cell(1, 0).text = "Laptop"
    table.cell(1, 1).text = "999"
    doc.save(path)

    text = llm.read_document_file(str(path))

    assert "# Contract heading" in text
    assert "Docx contract text" in text
    assert "| Item | Total |" in text
    assert "| Laptop | 999 |" in text


def test_xlsx_reader_extracts_sheet_cells(tmp_path):
    pytest.importorskip("anydoc")
    pytest.importorskip("openpyxl")
    import openpyxl

    from core.llm_clients import client as llm

    path = tmp_path / "sample.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Invoices"
    ws["A1"] = "Item"
    ws["B1"] = "Total"
    ws["A2"] = "Laptop"
    ws["B2"] = 999
    archive = wb.create_sheet("Archive")
    archive["A1"] = "Closed"
    wb.save(path)

    text = llm.read_document_file(str(path))

    assert "## Invoices" in text
    assert "## Archive" in text
    assert "| Item | Total |" in text
    assert "| Laptop | 999 |" in text


def test_pptx_reader_extracts_slide_text(tmp_path):
    pytest.importorskip("anydoc")
    pytest.importorskip("pptx")
    from pptx import Presentation

    from core.llm_clients import client as llm

    path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Deck contract title"
    slide.notes_slide.notes_text_frame.text = "Speaker-only detail"
    prs.save(path)

    text = llm.read_document_file(str(path))

    assert "## Deck contract title" in text
    assert "> Speaker-only detail" in text


def test_odt_reader_extracts_paragraph_text(tmp_path):
    pytest.importorskip("anydoc")
    pytest.importorskip("odf")
    from odf import text as odf_text
    from odf.opendocument import OpenDocumentText

    from core.llm_clients import client as llm

    path = tmp_path / "sample.odt"
    doc = OpenDocumentText()
    doc.text.addElement(odf_text.P(text="ODT contract text"))
    doc.save(str(path))

    assert "ODT contract text" in llm.read_document_file(str(path))


def test_anydoc_rtf_reader_preserves_inline_formatting(tmp_path):
    pytest.importorskip("anydoc")
    from core.llm_clients import client as llm

    path = tmp_path / "sample.rtf"
    path.write_text(r"{\rtf1\ansi\b Bold title\b0\par Plain body}", encoding="ascii")

    text = llm.read_document_file(str(path))

    assert "**Bold title**" in text
    assert "Plain body" in text


def test_anydoc_unavailable_uses_legacy_docx_fallback(tmp_path, monkeypatch):
    pytest.importorskip("docx")
    from docx import Document

    from core.llm_clients import documents

    path = tmp_path / "legacy.docx"
    doc = Document()
    doc.add_paragraph("Legacy fallback text")
    doc.save(path)
    monkeypatch.setattr(documents, "_load_anydoc", lambda: None)

    assert "Legacy fallback text" in documents._read_document_file(str(path))


def test_anydoc_conversion_failure_uses_legacy_fallback(tmp_path, monkeypatch):
    pytest.importorskip("docx")
    from docx import Document

    from core.llm_clients import documents

    class ConvertError(Exception):
        pass

    class ResourceLimitError(ConvertError):
        pass

    class EncryptedError(ConvertError):
        pass

    def fail_conversion(_path):
        raise ConvertError("new parser could not recover")

    parser = SimpleNamespace(
        ConvertError=ConvertError,
        ResourceLimitError=ResourceLimitError,
        EncryptedError=EncryptedError,
        to_markdown=fail_conversion,
    )
    path = tmp_path / "fallback.docx"
    doc = Document()
    doc.add_paragraph("Recovered by mature parser")
    doc.save(path)
    monkeypatch.setattr(documents, "_load_anydoc", lambda: parser)

    assert "Recovered by mature parser" in documents._read_document_file(str(path))


@pytest.mark.parametrize(
    ("error_name", "expected"),
    [
        ("ResourceLimitError", "parser safety limit exceeded"),
        ("EncryptedError", "encrypted or password-protected"),
    ],
)
def test_anydoc_security_failures_do_not_fall_back(tmp_path, monkeypatch, error_name, expected):
    from core.llm_clients import documents

    class ConvertError(Exception):
        pass

    class ResourceLimitError(ConvertError):
        pass

    class EncryptedError(ConvertError):
        pass

    error_type = {
        "ResourceLimitError": ResourceLimitError,
        "EncryptedError": EncryptedError,
    }[error_name]

    def fail_conversion(_path):
        raise error_type("blocked")

    parser = SimpleNamespace(
        ConvertError=ConvertError,
        ResourceLimitError=ResourceLimitError,
        EncryptedError=EncryptedError,
        to_markdown=fail_conversion,
    )
    path = tmp_path / "blocked.docx"
    path.write_bytes(b"not a real document")
    monkeypatch.setattr(documents, "_load_anydoc", lambda: parser)

    text = documents._read_document_file(str(path))

    assert text.startswith("Failed to read")
    assert expected in text


def test_anydoc_output_is_truncated_after_conversion(tmp_path, monkeypatch):
    from core.llm_clients import documents

    class ConvertError(Exception):
        pass

    parser = SimpleNamespace(
        ConvertError=ConvertError,
        ResourceLimitError=type("ResourceLimitError", (ConvertError,), {}),
        EncryptedError=type("EncryptedError", (ConvertError,), {}),
        to_markdown=lambda _path: "x" * 10_000,
    )
    path = tmp_path / "bounded.rtf"
    path.write_text(r"{\rtf1 bounded}", encoding="ascii")
    monkeypatch.setattr(documents, "_load_anydoc", lambda: parser)

    text = documents._read_document_file(str(path), max_chars=500)

    assert len(text) < 550
    assert "truncated" in text


def test_pdf_reader_dispatches_to_pdf_text_extractor(tmp_path, monkeypatch):
    from core.llm_clients import client as llm
    from core.llm_clients import documents

    path = tmp_path / "sample.pdf"
    path.write_bytes(b"%PDF-1.4\n% test placeholder\n")
    # PDF extraction lives in core.llm_clients.documents (re-exported via client);
    # patch it at the seam where read_document_file resolves it.
    monkeypatch.setattr(documents, "_read_pdf_text", lambda p, _max_chars: f"PDF text from {Path(p).name}")

    assert "PDF text from sample.pdf" in llm.read_document_file(str(path))
